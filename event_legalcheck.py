# -*- coding: utf-8 -*-
"""
event_legalcheck_v2_3_multisource_pptx_preview.py
------------------------------------------------------------
✅ v2.3 (멀티 입력 + PPTX 사전분석/컨펌 + 슬라이드 구조화 + 이미지객체 OCR)
- 입력 소스: URL / 파일 업로드(이미지, PPTX) / 텍스트 붙여넣기
- API Key 1개 (OCR/Gemini 선택 실행)
- PPTX: 업로드 시 먼저 분석(/analyze_pptx) → OCR 예상 호출 수 안내 → 사용자 컨펌 후 OCR 실행
- PPTX: python-pptx 기반 슬라이드별 텍스트 추출 + 슬라이드 내 이미지 객체 추출(OCR 선택)
- PPT 업로드 시: pptx로 저장 후 재업로드 안내
- 개발자 모드: Gemini raw JSON / raw text 토글
------------------------------------------------------------
"""

import os
import re
import time
import base64
import json
import hashlib
import webbrowser
from datetime import datetime
from io import BytesIO

import requests
from flask import Flask, request, render_template_string, jsonify, send_from_directory
from PIL import Image

# Playwright (sync)
from playwright.sync_api import sync_playwright, TimeoutError as PWTimeoutError

# PPTX text + images
try:
    from pptx import Presentation
except Exception:
    Presentation = None

# ------------------------------------------------------------
# 기본 설정
# ------------------------------------------------------------
OUTPUT_DIR = os.path.join(os.getcwd(), "outputs")
UPLOAD_DIR = os.path.join(os.getcwd(), "uploads")
TMP_DIR = os.path.join(os.getcwd(), "tmp_pptx")

os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(TMP_DIR, exist_ok=True)

VISION_ENDPOINT = "https://vision.googleapis.com/v1/images:annotate"
GEMINI_ENDPOINT = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash-exp:generateContent"

ALLOWED_IMAGE_EXT = {".png", ".jpg", ".jpeg", ".webp"}
ALLOWED_PPTX_EXT = {".pptx"}
ALLOWED_PPT_EXT = {".ppt"}  # 안내용

# in-memory file token storage (simple)
PPTX_TOKENS = {}  # token -> {"path":..., "filename":..., "created":..., "analyze": {...}}


# ------------------------------------------------------------
# Utility
# ------------------------------------------------------------
def normalize_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def escape_html(text):
    if text is None:
        return ""
    return (text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&lt;").replace(">", "&gt;"))


def normalize_api_key(k: str) -> str:
    k = (k or "").strip().replace("\n", "").replace("\r", "")
    k = k.replace(" ", "")
    k = k.replace("key=", "")
    k = k.replace('"', "").replace("'", "")
    return k


def safe_filename(name: str) -> str:
    name = name or "upload"
    name = re.sub(r"[^\w\.\-]", "_", name)
    return name[:140]


def get_ext(filename: str) -> str:
    _, ext = os.path.splitext(filename or "")
    return (ext or "").lower().strip()


def sha1_bytes(b: bytes) -> str:
    h = hashlib.sha1()
    h.update(b)
    return h.hexdigest()


# ------------------------------------------------------------
# OCR (Google Vision REST + API KEY)
# ------------------------------------------------------------
def _pil_to_base64(img: Image.Image) -> str:
    buf = BytesIO()
    img.save(buf, format="PNG")
    return base64.b64encode(buf.getvalue()).decode("utf-8")


def ocr_google_vision(img: Image.Image, api_key: str, use_document_text=True) -> str:
    api_key = normalize_api_key(api_key)
    if not api_key:
        raise RuntimeError("API KEY가 비어있습니다.")

    feature = "DOCUMENT_TEXT_DETECTION" if use_document_text else "TEXT_DETECTION"
    payload = {
        "requests": [
            {
                "image": {"content": _pil_to_base64(img)},
                "features": [{"type": feature}],
                "imageContext": {"languageHints": ["ko", "en"]}
            }
        ]
    }

    url = f"{VISION_ENDPOINT}?key={api_key}"
    resp = requests.post(url, json=payload, timeout=60)
    if resp.status_code != 200:
        raise RuntimeError(f"Vision REST 호출 실패: HTTP {resp.status_code} - {resp.text[:500]}")

    data = resp.json()
    if "error" in data:
        raise RuntimeError(f"Vision API error: {data['error']}")

    r0 = (data.get("responses") or [{}])[0]
    if "error" in r0:
        raise RuntimeError(f"Vision API error: {r0['error']}")

    if use_document_text:
        return ((r0.get("fullTextAnnotation") or {}).get("text")) or ""

    anns = r0.get("textAnnotations") or []
    return anns[0].get("description", "") if anns else ""


# ------------------------------------------------------------
# Gemini API (REST)
# ------------------------------------------------------------
def review_with_gemini_rest(
        source_label: str,
        visible_text: str,
        hidden_text: str,
        iframe_text: str,
        ocr_text: str,
        rule_results: list,
        findings: list,
        page_type: str,
        api_key: str,
        slide_context: str = ""
) -> dict:

    api_key = normalize_api_key(api_key)
    if not api_key:
        return {
            "gemini_review": "API Key가 제공되지 않았습니다.",
            "confidence_score": 0,
            "additional_findings": [],
            "recommendations": [],
            "raw_json": None,
            "raw_text": ""
        }

    response_text = ""
    try:
        prompt = f"""당신은 한국의 법률 및 규정 준수 전문가입니다. 다음 검토 대상에 대한 법적 검토 결과를 최종 검증해주세요.

**검토 대상**: {source_label}
**페이지 유형**: {page_type}

{slide_context}

**수집된 텍스트 정보**:
1. 메인 영역 텍스트 (처음 3000자):
{visible_text[:3000]}

2. 숨김 레이어 텍스트:
{hidden_text[:1500] if hidden_text else "(없음)"}

3. iframe 텍스트:
{iframe_text[:1000] if iframe_text else "(없음)"}

4. OCR 텍스트:
{ocr_text[:1500] if ocr_text else "(없음)"}

**규칙 기반 검사 결과**:
{json.dumps(rule_results[:10], ensure_ascii=False, indent=2)}

**발견된 위반 사항**:
{json.dumps(findings, ensure_ascii=False, indent=2) if findings else "(위반 사항 없음)"}

**검토 요청 사항**:
1. 규칙 기반 검사 결과가 정확한지 검증해주세요
2. 규칙에서 놓친 법적 이슈가 있는지 확인해주세요
3. 특히 다음 항목을 중점적으로 검토해주세요:
   - 개인정보보호법 (마케팅 동의, 제3자 제공 동의)
   - 표시광고법 (경품 표시, 사은품 조건)
   - 전자상거래법 (환불/교환 정책)
   - 약관규제법 (불공정 조항)

4. 응답은 반드시 다음 JSON 형식으로만 작성해주세요 (다른 텍스트 없이).
   - 가능하면 '추가발견사항' 항목에 슬라이드 번호가 특정되는 경우 "슬라이드 3"처럼 포함해주세요.
{{
  "종합검토의견": "전반적인 법적 준수 상태 평가 (200자 이내)",
  "신뢰도점수": 85,
  "추가발견사항": [
    {{"심각도": "CRITICAL", "항목": "항목명", "설명": "설명", "근거": "근거 텍스트"}}
  ],
  "개선권고사항": ["권고1", "권고2"]
}}
"""

        gemini_url = f"{GEMINI_ENDPOINT}?key={api_key}"

        payload = {
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {
                "temperature": 0.3,
                "topK": 40,
                "topP": 0.95,
                "maxOutputTokens": 2048,
            }
        }

        resp = requests.post(gemini_url, json=payload, timeout=60)
        if resp.status_code != 200:
            raise RuntimeError(f"Gemini API 호출 실패: HTTP {resp.status_code} - {resp.text[:500]}")

        data = resp.json()

        candidates = data.get("candidates", [])
        if not candidates:
            raise RuntimeError("Gemini API 응답에 candidates가 없습니다")

        content = candidates[0].get("content", {})
        parts = content.get("parts", [])
        if not parts:
            raise RuntimeError("Gemini API 응답에 parts가 없습니다")

        response_text = (parts[0].get("text", "") or "").strip()

        extracted = response_text
        if "```json" in extracted:
            extracted = extracted.split("```json", 1)[1]
            extracted = extracted.split("```", 1)[0].strip()
        elif "```" in extracted:
            extracted = extracted.split("```", 1)[1]
            extracted = extracted.split("```", 1)[0].strip()

        gemini_result = json.loads(extracted)

        return {
            "gemini_review": gemini_result.get("종합검토의견", ""),
            "confidence_score": int(gemini_result.get("신뢰도점수", 0)),
            "additional_findings": gemini_result.get("추가발견사항", []),
            "recommendations": gemini_result.get("개선권고사항", []),
            "raw_json": gemini_result,
            "raw_text": response_text
        }

    except json.JSONDecodeError as e:
        return {
            "gemini_review": f"Gemini API 응답 파싱 실패: {str(e)}",
            "confidence_score": 0,
            "additional_findings": [],
            "recommendations": [],
            "raw_json": None,
            "raw_text": response_text[:8000]
        }
    except Exception as e:
        return {
            "gemini_review": f"Gemini API 호출 중 오류: {str(e)}",
            "confidence_score": 0,
            "additional_findings": [],
            "recommendations": [],
            "raw_json": None,
            "raw_text": response_text[:8000]
        }


# ------------------------------------------------------------
# Page type classification
# ------------------------------------------------------------
PAGE_TYPES = {
    "OUTBOUND_EVENT": [r"아웃바운드", r"전화\s*상담", r"상담\s*전화", r"TM", r"콜", r"콜백"],
    "EVENT": [r"이벤트", r"경품", r"추첨", r"당첨", r"쿠폰"],
    "INSURANCE": [r"보험", r"설계사", r"보장", r"가입"],
    "ECOMMERCE": [r"결제", r"주문", r"배송", r"장바구니"],
    "FORM_ONLY": [r"이름", r"연락처", r"휴대폰", r"상담", r"신청"],
}


def classify_page_type(all_text: str):
    hits = []
    for t, pats in PAGE_TYPES.items():
        for p in pats:
            if re.search(p, all_text or "", re.I):
                hits.append(t)
                break

    if not hits:
        return "GENERIC", ["유형 키워드 미검출 → GENERIC로 분류"]

    priority = ["OUTBOUND_EVENT", "EVENT", "INSURANCE", "ECOMMERCE", "FORM_ONLY"]
    hits_sorted = [h for h in priority if h in hits]
    picked = hits_sorted[0]
    reason = [f"{picked} 관련 키워드 검출 → {picked}로 분류"]
    return picked, reason


# ------------------------------------------------------------
# Playwright 렌더링 (URL)
# ------------------------------------------------------------
def render_page_text_playwright(url: str, timeout_ms: int = 20000):
    screenshot_path = None
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(
            viewport={"width": 1280, "height": 720},
            user_agent="Mozilla/5.0 (LegalLandingChecker/v2.3)"
        )
        page = context.new_page()

        try:
            page.goto(url, wait_until="networkidle", timeout=timeout_ms)
        except PWTimeoutError:
            pass

        # 스크롤
        try:
            page.evaluate("""
                          () => new Promise((resolve) => {
                              let total = 0;
                              const distance = 700;
                              const timer = setInterval(() => {
                                  window.scrollBy(0, distance);
                                  total += distance;
                                  if (total >= document.body.scrollHeight) {
                                      clearInterval(timer);
                                      window.scrollTo(0, 0);
                                      resolve(true);
                                  }
                              }, 120);
                          })
                          """)
        except Exception:
            pass

        visible_text = ""
        try:
            visible_text = page.evaluate("""
                                         () => {
                                             const isVisible = (el) => {
                                                 if (!el) return false;
                                                 const style = window.getComputedStyle(el);
                                                 if (!style) return false;
                                                 if (style.display === 'none' || style.visibility === 'hidden' || style.opacity === '0') return false;
                                                 const rect = el.getBoundingClientRect();
                                                 if (rect.width === 0 || rect.height === 0) return false;
                                                 return true;
                                             };
                                             const walker = document.createTreeWalker(document.body, NodeFilter.SHOW_ELEMENT, null);
                                             const texts = [];
                                             let node = walker.currentNode;
                                             while (node) {
                                                 if (isVisible(node)) {
                                                     const t = (node.innerText || '').trim();
                                                     if (t && t.length >= 2) texts.push(t);
                                                 }
                                                 node = walker.nextNode();
                                             }
                                             return Array.from(new Set(texts)).join("\\n");
                                         }
                                         """)
        except Exception:
            pass

        hidden_text = ""
        try:
            hidden_text = page.evaluate("""
                                        () => {
                                            const isHidden = (el) => {
                                                if (!el) return false;
                                                if (el.hidden) return true;
                                                if (el.getAttribute && el.getAttribute('aria-hidden') === 'true') return true;
                                                const style = window.getComputedStyle(el);
                                                if (!style) return false;
                                                if (style.display === 'none' || style.visibility === 'hidden' || style.opacity === '0') return true;
                                                const rect = el.getBoundingClientRect();
                                                if (rect.width === 0 || rect.height === 0) return true;
                                                return false;
                                            };
                                            const walker = document.createTreeWalker(document.body, NodeFilter.SHOW_ELEMENT, null);
                                            const texts = [];
                                            let node = walker.currentNode;
                                            while (node) {
                                                if (isHidden(node)) {
                                                    const t = (node.innerText || '').trim();
                                                    if (t && t.length >= 2) texts.push(t);
                                                }
                                                node = walker.nextNode();
                                            }
                                            return Array.from(new Set(texts)).join("\\n");
                                        }
                                        """)
        except Exception:
            pass

        iframe_texts = []
        try:
            for frame in page.frames:
                if frame == page.main_frame:
                    continue
                try:
                    t = frame.evaluate("() => document.body ? (document.body.innerText || '') : ''")
                    if t:
                        iframe_texts.append(t.strip())
                except Exception:
                    continue
        except Exception:
            pass

        iframe_text = "\n\n".join(iframe_texts)

        try:
            screenshot_path = os.path.join(OUTPUT_DIR, f"shot_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
            page.screenshot(path=screenshot_path, full_page=True)
        except Exception:
            screenshot_path = None

        context.close()
        browser.close()

    return normalize_text(visible_text), normalize_text(hidden_text), normalize_text(iframe_text), screenshot_path


# ------------------------------------------------------------
# PPTX 분석 / 추출
# ------------------------------------------------------------
def pptx_analyze(filepath: str) -> dict:
    """PPTX 사전 분석: 슬라이드 수, 이미지 객체 수, 텍스트 존재 여부"""
    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되어 있지 않습니다. `pip install python-pptx` 후 재시도하세요.")

    prs = Presentation(filepath)
    slide_count = len(prs.slides)
    image_count = 0
    has_text_layer = False

    for slide in prs.slides:
        for shape in slide.shapes:
            # text check
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                has_text_layer = True
            # image check
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE == 13
                image_count += 1

    return {
        "slide_count": slide_count,
        "image_count": image_count,
        "has_text_layer": has_text_layer,
        "estimated_ocr_calls": image_count
    }


def pptx_extract_slides(filepath: str) -> list:
    """
    슬라이드별 텍스트 + 이미지 blob 추출
    returns:
      [
        {"slide_no":1, "text":"...", "images":[{"path":"...", "sha1":"..."}]},
        ...
      ]
    """
    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되어 있지 않습니다. `pip install python-pptx` 후 재시도하세요.")

    prs = Presentation(filepath)
    slides_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        images = []

        for shape in slide.shapes:
            # text
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    slide_texts.append(t)

            # image
            if shape.shape_type == 13:  # picture
                try:
                    img_blob = shape.image.blob
                    img_sha1 = sha1_bytes(img_blob)
                    img_path = os.path.join(TMP_DIR, f"pptx_slide{idx}_img_{img_sha1}.png")
                    if not os.path.exists(img_path):
                        with open(img_path, "wb") as fw:
                            fw.write(img_blob)
                    images.append({"path": img_path, "sha1": img_sha1})
                except Exception:
                    continue

        slides_data.append({
            "slide_no": idx,
            "text": normalize_text("\n".join(slide_texts)),
            "images": images,
            "ocr": ""
        })

    return slides_data


def pptx_run_image_ocr(slides_data: list, api_key: str) -> list:
    """슬라이드별 이미지 객체만 OCR 수행"""
    for s in slides_data:
        ocr_chunks = []
        for img in s.get("images", []):
            try:
                im = Image.open(img["path"])
                t = ocr_google_vision(im, api_key)
                t = normalize_text(t)
                if t:
                    ocr_chunks.append(t)
            except Exception:
                continue
        s["ocr"] = normalize_text("\n\n".join(ocr_chunks))
    return slides_data


def pptx_merge_text(slides_data: list) -> (str, str, str):
    """
    슬라이드 구조를 전체 텍스트로 합치기
    returns:
      merged_visible_text, merged_ocr_text, slide_context_summary
    """
    merged_texts = []
    merged_ocr = []
    slide_context_lines = []

    for s in slides_data:
        no = s["slide_no"]
        txt = s.get("text") or ""
        ocr = s.get("ocr") or ""

        if txt:
            merged_texts.append(f"[슬라이드 {no}]\n{txt}")
        if ocr:
            merged_ocr.append(f"[슬라이드 {no} OCR]\n{ocr}")

        slide_context_lines.append(
            f"- 슬라이드 {no}: 텍스트 {len(txt)}자 / OCR {len(ocr)}자 / 이미지 {len(s.get('images', []))}개"
        )

    slide_context_summary = "※ 슬라이드별 구조 정보:\n" + "\n".join(slide_context_lines)
    return normalize_text("\n\n".join(merged_texts)), normalize_text("\n\n".join(merged_ocr)), slide_context_summary


# ------------------------------------------------------------
# Rules (간소화 — 프로님 기존 full run_rules로 교체 가능)
# ------------------------------------------------------------
RULES = [
    (
        "PIPA_MARKETING_REQUIRED",
        "CRITICAL",
        "개인정보보호법 제22조(동의) - 마케팅 수신 동의는 선택 동의 원칙",
        "마케팅 정보 수신 동의가 '필수'로 설계될 경우 리스크가 큽니다.",
        [r"마케팅\s*정보\s*수신.*필수", r"마케팅.*수신.*동의.*필수"],
        20,
        "PIPA",
        {"applicability": "ALWAYS"}
    ),
]


def run_rules(visible_text, hidden_text="", iframe_text="", ocr_text="", enable_hidden=True, enable_ocr=False,
              page_type="GENERIC"):
    base_score = 100
    findings = []
    rule_audit = []
    passed_items = []
    na_items = []
    satisfied_in_hidden = []

    all_text = "\n".join([visible_text or "", hidden_text or "", iframe_text or "", ocr_text or ""])
    all_text = normalize_text(all_text)

    for rule in RULES:
        rule_id, severity, title, desc, patterns, penalty, category, meta = rule
        found = False
        for p in patterns:
            if re.search(p, all_text, re.I):
                found = True
                break

        if found:
            base_score -= penalty
            findings.append({
                "severity": severity,
                "title": title,
                "desc": desc
            })
            rule_audit.append({
                "id": rule_id,
                "title": title,
                "status": "FAIL",
                "reason": "패턴 매칭"
            })
        else:
            rule_audit.append({
                "id": rule_id,
                "title": title,
                "status": "PASS",
                "reason": "미검출"
            })

    return max(base_score, 0), findings, satisfied_in_hidden, rule_audit, passed_items, na_items, all_text


# ------------------------------------------------------------
# Flask App
# ------------------------------------------------------------
app = Flask(__name__)


@app.route("/")
def index():
    return render_template_string(r"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Legal Landing Checker v2.3</title>
    <style>
        body { font-family: 'Segoe UI', sans-serif; max-width: 1020px; margin: 40px auto; padding: 20px; }
        h1 { color: #1a73e8; }
        .form-group { margin: 18px 0; }
        label { display: block; margin-bottom: 6px; font-weight: 600; }
        input[type="text"], textarea { width: 100%; padding: 10px; border: 1px solid #ddd; border-radius: 4px; box-sizing: border-box; }
        textarea { min-height: 110px; font-family: monospace; }
        button { background: #1a73e8; color: white; padding: 12px 22px; border: none; border-radius: 4px; cursor: pointer; font-size: 15px; }
        button:hover { background: #1557b0; }
        button.secondary { background: #444; }
        button.secondary:hover { background: #222; }
        .options { background: #f8f9fa; padding: 15px; border-radius: 6px; margin: 18px 0; }
        .checkbox-group { margin: 10px 0; }
        .checkbox-group label { display: inline-block; margin-right: 14px; font-weight: normal; }
        #status { margin-top: 18px; padding: 15px; border-radius: 4px; display: none; }
        .success { background: #d4edda; color: #155724; }
        .error { background: #f8d7da; color: #721c24; }
        .warn { background:#fff3cd; color:#856404; }
        .new-badge { background: #ea4335; color: white; padding: 2px 6px; border-radius: 3px; font-size: 11px; font-weight: bold; margin-left: 8px; }
        .advanced-box { margin-top: 12px; padding: 12px; background: #fff; border: 1px dashed #ccc; border-radius: 6px; display:none; }
        .advanced-title { font-weight: 700; color:#333; margin-bottom: 8px; }
        .tab-box { display: flex; gap: 10px; margin-top: 8px; }
        .tab-box label { font-weight: normal; }
        .hint { color:#666; font-size: 13px; line-height: 1.4; }
        .panel { display:none; margin-top: 12px; }
        .panel.active { display:block; }
        .small { font-size: 12px; color:#777; }
        .preview-box { margin-top: 12px; padding: 12px; background:#fff; border:1px solid #ddd; border-radius: 8px; display:none; }
        .preview-row { display:flex; justify-content:space-between; padding: 6px 0; border-bottom:1px dashed #eee; }
        .preview-row:last-child { border-bottom:none; }
        .preview-actions { display:flex; gap: 10px; margin-top: 12px; }
        code { background:#f1f3f4; padding:2px 6px; border-radius:4px; }
    </style>
</head>
<body>
    <h1>🔍 Legal Landing Checker v2.3</h1>
    <p><strong>PPTX 고도화</strong> - 슬라이드 구조화 + 이미지 객체 OCR(컨펌 후 실행)</p>

    <div class="options">
        <h3>입력 방식 선택</h3>
        <div class="tab-box">
            <label><input type="radio" name="source_type" value="URL" checked /> URL</label>
            <label><input type="radio" name="source_type" value="FILE" /> 파일 첨부</label>
            <label><input type="radio" name="source_type" value="TEXT" /> 텍스트 붙여넣기</label>
        </div>

        <div id="panel_url" class="panel active">
            <div class="form-group">
                <label>검사할 URL:</label>
                <input type="text" id="url" placeholder="https://example.com/landing" />
            </div>
        </div>

        <div id="panel_file" class="panel">
            <div class="form-group">
                <label>검수할 파일 업로드 (이미지/PPTX)</label>
                <input type="file" id="file_input" />
                <div class="hint" style="margin-top:8px;">
                    • 지원: PNG/JPG/WEBP, PPTX<br/>
                    • PPT(.ppt)는 구버전 포맷으로 자동 추출이 어려워 <b>PPTX로 저장 후 업로드</b> 권장<br/>
                    • PPTX의 OCR(이미지 객체 OCR)은 <b>업로드 후 사전 분석 결과를 보고 컨펌</b>해야 실행됩니다.
                </div>

                <div id="pptx_preview" class="preview-box">
                    <h4 style="margin:0 0 8px 0;">📌 PPTX 분석 결과 (OCR 실행 전)</h4>
                    <div class="preview-row"><span>파일</span><span id="pv_name"></span></div>
                    <div class="preview-row"><span>슬라이드 수</span><span id="pv_slides"></span></div>
                    <div class="preview-row"><span>이미지 객체 수</span><span id="pv_imgs"></span></div>
                    <div class="preview-row"><span>텍스트 레이어 존재</span><span id="pv_textlayer"></span></div>
                    <div class="preview-row"><span>예상 OCR 호출 수</span><span id="pv_calls"></span></div>

                    <div class="preview-actions">
                        <button onclick="runCheckWithToken(false)" class="secondary">OCR 없이 검사 실행</button>
                        <button onclick="runCheckWithToken(true)">OCR 포함 검사 실행</button>
                    </div>
                    <div class="small" style="margin-top:10px;">
                        ※ OCR 포함 실행은 Vision API가 활성화된 Key가 필요합니다.
                    </div>
                </div>

            </div>
        </div>

        <div id="panel_text" class="panel">
            <div class="form-group">
                <label>검수할 텍스트 붙여넣기</label>
                <textarea id="input_text" placeholder="여기에 검수할 텍스트를 붙여넣으세요."></textarea>
                <div class="small">※ 텍스트 모드에서는 OCR은 적용되지 않습니다.</div>
            </div>
        </div>
    </div>

    <div class="options">
        <h3>검사 옵션</h3>
        <div class="checkbox-group">
            <label><input type="checkbox" id="enable_rendering" checked /> 렌더링 모드 (URL 전용)</label>
            <label><input type="checkbox" id="enable_hidden" checked /> 숨김 레이어 검사 (URL 전용)</label>
        </div>
        <div class="checkbox-group">
            <label><input type="checkbox" id="enable_ocr" /> OCR 검사 (이미지/URL 스크린샷)</label>
            <label>
                <input type="checkbox" id="enable_gemini" checked /> 
                AI 최종 검토 (Gemini API)
                <span class="new-badge">NEW</span>
            </label>
        </div>

        <div class="checkbox-group">
            <label><input type="checkbox" id="dev_mode" /> 개발자 모드 (고급 옵션)</label>
        </div>

        <div id="advanced_box" class="advanced-box">
            <div class="advanced-title">고급 옵션</div>
            <div class="hint">
                • Gemini 원문 JSON/원문 응답 노출 (리포트 하단에 토글 표시)<br/>
                • 디버깅/품질 점검 목적으로만 사용 권장
            </div>
        </div>

        <div class="form-group">
            <label>Google Cloud API Key: <small>(OCR/Gemini 선택 사용 — 동일 Key 사용 가능)</small></label>
            <textarea id="api_key" placeholder="AIzaSy..."></textarea>
            <small id="key_hint" class="hint"></small>
        </div>
    </div>

    <button onclick="mainRun()">검사 실행</button>

    <div id="status"></div>

    <script>
        let pptxFileToken = null;
        let pptxFileName = "";

        function getSourceType() {
            const r = document.querySelector('input[name="source_type"]:checked');
            return r ? r.value : "URL";
        }

        function setPanels() {
            const t = getSourceType();
            document.getElementById("panel_url").classList.toggle("active", t === "URL");
            document.getElementById("panel_file").classList.toggle("active", t === "FILE");
            document.getElementById("panel_text").classList.toggle("active", t === "TEXT");

            // URL 전용 옵션 disable 처리
            const enable_rendering = document.getElementById("enable_rendering");
            const enable_hidden = document.getElementById("enable_hidden");

            if (t !== "URL") {
                enable_rendering.checked = false;
                enable_hidden.checked = false;
                enable_rendering.disabled = true;
                enable_hidden.disabled = true;
            } else {
                enable_rendering.disabled = false;
                enable_hidden.disabled = false;
            }

            // TEXT 모드 OCR 비활성화
            const enable_ocr = document.getElementById("enable_ocr");
            if (t === "TEXT") {
                enable_ocr.checked = false;
                enable_ocr.disabled = true;
            } else {
                enable_ocr.disabled = false;
            }

            updateKeyHint();
        }

        function updateKeyHint() {
            const enable_ocr = document.getElementById('enable_ocr').checked;
            const enable_gemini = document.getElementById('enable_gemini').checked;
            const hint = document.getElementById('key_hint');

            if (!enable_ocr && !enable_gemini) {
                hint.innerHTML = `ℹ️ 현재는 <b>규칙 기반 검사만 실행</b>됩니다. (API Key 입력 불필요)`;
            } else if (enable_ocr && !enable_gemini) {
                hint.innerHTML = `
                    ℹ️ <b>OCR 검사만 실행</b>됩니다.<br/>
                    • API Key 입력 필요<br/>
                    • <b>Vision API 활성화 필요</b>
                `;
            } else if (!enable_ocr && enable_gemini) {
                hint.innerHTML = `
                    ℹ️ <b>Gemini AI 검토만 실행</b>됩니다.<br/>
                    • API Key 입력 필요<br/>
                    • <b>Generative Language API 활성화 필요</b><br/>
                    • Vision API는 활성화하지 않아도 됩니다.
                `;
            } else {
                hint.innerHTML = `
                    ℹ️ <b>OCR + Gemini AI 검토가 모두 실행</b>됩니다.<br/>
                    • API Key 입력 필요<br/>
                    • <b>Vision API + Generative Language API 둘 다 활성화 필요</b>
                `;
            }
        }

        function toggleAdvanced() {
            const dev = document.getElementById('dev_mode').checked;
            const box = document.getElementById('advanced_box');
            box.style.display = dev ? 'block' : 'none';
        }

        function showStatus(cls, msg) {
            const status = document.getElementById('status');
            status.style.display = 'block';
            status.className = cls;
            status.innerHTML = msg;
        }

        async function analyzePptxIfNeeded(file) {
            // file이 pptx면 /analyze_pptx 호출 후 preview 표시
            const ext = (file.name || "").toLowerCase().split(".").pop();
            if (ext !== "pptx") return false;

            const api_key = document.getElementById('api_key').value.trim();

            showStatus("", "PPTX 분석 중... (OCR 예상 호출 수 산정)");
            const formData = new FormData();
            formData.append("file", file);

            const resp = await fetch("/analyze_pptx", { method:"POST", body: formData });
            const result = await resp.json();

            if (result.error) {
                showStatus("error", "오류: " + result.error);
                return false;
            }

            pptxFileToken = result.file_token;
            pptxFileName = result.filename || file.name;

            // preview 표시
            document.getElementById("pptx_preview").style.display = "block";
            document.getElementById("pv_name").textContent = pptxFileName;
            document.getElementById("pv_slides").textContent = result.slide_count;
            document.getElementById("pv_imgs").textContent = result.image_count;
            document.getElementById("pv_textlayer").textContent = result.has_text_layer ? "있음" : "없음";
            document.getElementById("pv_calls").textContent = result.estimated_ocr_calls;

            showStatus("warn", "📌 PPTX 분석 완료. 아래에서 OCR 포함/제외를 선택해 검사 실행하세요.");
            return true;
        }

        async function mainRun() {
            const source_type = getSourceType();
            if (source_type === "FILE") {
                const f = document.getElementById("file_input").files[0];
                if (!f) { alert("파일을 첨부하세요"); return; }
                const ext = (f.name || "").toLowerCase().split(".").pop();

                // PPTX면 분석 단계로 유도
                if (ext === "pptx") {
                    await analyzePptxIfNeeded(f);
                    return;
                }
            }

            // PPTX 아닌 경우는 바로 runCheck
            await runCheckDirect();
        }

        async function runCheckWithToken(doOcrPptx) {
            // pptx preview 이후 token 기반 실행
            if (!pptxFileToken) {
                alert("PPTX 토큰이 없습니다. 파일을 다시 업로드해주세요.");
                return;
            }
            const api_key = document.getElementById('api_key').value.trim();
            const enable_gemini = document.getElementById('enable_gemini').checked;

            // OCR 포함이면 api_key 필수
            if (doOcrPptx && !api_key) {
                alert("PPTX OCR 포함 실행을 위해 API Key가 필요합니다.");
                return;
            }
            if (enable_gemini && !api_key) {
                alert("Gemini 검토를 위해 API Key가 필요합니다.");
                return;
            }

            showStatus("", "검사 진행 중... (PPTX 토큰 기반)");

            const formData = new FormData();
            formData.append("source_type", "PPTX_TOKEN");
            formData.append("pptx_token", pptxFileToken);
            formData.append("pptx_ocr", doOcrPptx ? "1" : "0");

            // 옵션
            formData.append("enable_ocr", "0"); // pptx는 별도 플래그로 처리
            formData.append("enable_gemini", enable_gemini ? "1" : "0");
            formData.append("dev_mode", document.getElementById('dev_mode').checked ? "1" : "0");
            formData.append("api_key", api_key);

            const resp = await fetch('/check', { method:"POST", body: formData });
            const result = await resp.json();

            if (result.error) {
                showStatus("error", "오류: " + result.error);
            } else {
                showStatus("success", '✅ 검사 완료! <a href="/report/' + result.report + '" target="_blank" style="font-weight:bold; color:#155724;">결과 보기</a>');
            }
        }

        async function runCheckDirect() {
            const source_type = getSourceType();
            const enable_ocr = document.getElementById('enable_ocr').checked;
            const enable_gemini = document.getElementById('enable_gemini').checked;
            const api_key = document.getElementById('api_key').value.trim();

            if ((enable_ocr || enable_gemini) && !api_key) {
                alert("OCR 또는 Gemini를 사용하려면 API Key 입력이 필요합니다.");
                return;
            }

            showStatus("", "검사 진행 중...");

            const formData = new FormData();
            formData.append("source_type", source_type);
            formData.append("enable_rendering", document.getElementById('enable_rendering').checked ? "1" : "0");
            formData.append("enable_hidden", document.getElementById('enable_hidden').checked ? "1" : "0");
            formData.append("enable_ocr", enable_ocr ? "1" : "0");
            formData.append("enable_gemini", enable_gemini ? "1" : "0");
            formData.append("dev_mode", document.getElementById('dev_mode').checked ? "1" : "0");
            formData.append("api_key", api_key);

            if (source_type === "URL") {
                const url = document.getElementById('url').value.trim();
                if (!url) { alert("URL을 입력하세요"); return; }
                formData.append("url", url);
            } else if (source_type === "TEXT") {
                const input_text = document.getElementById('input_text').value.trim();
                if (!input_text) { alert("텍스트를 입력/붙여넣기 하세요"); return; }
                formData.append("input_text", input_text);
            } else if (source_type === "FILE") {
                const f = document.getElementById('file_input').files[0];
                if (!f) { alert("파일을 첨부하세요"); return; }
                formData.append("file", f);
            }

            const resp = await fetch('/check', { method:'POST', body: formData });
            const result = await resp.json();

            if (result.error) {
                showStatus("error", "오류: " + result.error);
            } else {
                showStatus("success", '✅ 검사 완료! <a href="/report/' + result.report + '" target="_blank" style="font-weight:bold; color:#155724;">결과 보기</a>');
            }
        }

        document.addEventListener("DOMContentLoaded", () => {
            document.querySelectorAll('input[name="source_type"]').forEach(el => el.addEventListener('change', setPanels));
            document.getElementById('enable_ocr').addEventListener('change', updateKeyHint);
            document.getElementById('enable_gemini').addEventListener('change', updateKeyHint);
            document.getElementById('dev_mode').addEventListener('change', toggleAdvanced);
            setPanels();
            updateKeyHint();
            toggleAdvanced();

            // 파일 선택 시 pptx면 자동 분석
            document.getElementById("file_input").addEventListener("change", async (e) => {
                const f = e.target.files[0];
                if (!f) return;
                const ext = (f.name || "").toLowerCase().split(".").pop();
                if (ext === "pptx") {
                    await analyzePptxIfNeeded(f);
                } else {
                    document.getElementById("pptx_preview").style.display = "none";
                }
            });
        });
    </script>
</body>
</html>
    """)


# ------------------------------------------------------------
# PPTX 분석 API
# ------------------------------------------------------------
@app.route("/analyze_pptx", methods=["POST"])
def analyze_pptx():
    try:
        if "file" not in request.files:
            return jsonify({"error": "PPTX 파일이 첨부되지 않았습니다"}), 400

        f = request.files["file"]
        if not f or not f.filename:
            return jsonify({"error": "파일명이 유효하지 않습니다"}), 400

        filename = safe_filename(f.filename)
        ext = get_ext(filename)

        if ext in ALLOWED_PPT_EXT:
            return jsonify({"error": "PPT(.ppt)는 구버전 포맷입니다. PPTX로 저장 후 업로드해주세요."}), 400

        if ext not in ALLOWED_PPTX_EXT:
            return jsonify({"error": "PPTX 파일만 분석 가능합니다."}), 400

        save_path = os.path.join(UPLOAD_DIR, f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{filename}")
        f.save(save_path)

        a = pptx_analyze(save_path)
        token = f"pptx_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{hashlib.md5(save_path.encode()).hexdigest()[:10]}"
        PPTX_TOKENS[token] = {
            "path": save_path,
            "filename": filename,
            "created": time.time(),
            "analyze": a
        }

        return jsonify({
            "ok": True,
            "file_token": token,
            "filename": filename,
            **a
        })

    except Exception as e:
        import traceback
        return jsonify({"error": f"{str(e)}\n\n{traceback.format_exc()}"}), 500


# ------------------------------------------------------------
# /check : multipart 기반 처리 (URL/FILE/TEXT/PPTX_TOKEN)
# ------------------------------------------------------------
@app.route("/check", methods=["POST"])
def check():
    try:
        source_type = (request.form.get("source_type") or "URL").upper().strip()

        enable_rendering = request.form.get("enable_rendering") == "1"
        enable_hidden = request.form.get("enable_hidden") == "1"
        enable_ocr = request.form.get("enable_ocr") == "1"
        enable_gemini = request.form.get("enable_gemini") == "1"
        dev_mode = request.form.get("dev_mode") == "1"
        api_key = request.form.get("api_key", "")

        source_label = ""
        visible_text, hidden_text, iframe_text, screenshot_path = "", "", "", None
        ocr_text = ""
        slides_data = None
        slide_context_summary = ""

        # -----------------------------
        # 1) 입력 소스별 텍스트 확보
        # -----------------------------
        if source_type == "URL":
            url = (request.form.get("url") or "").strip()
            if not url:
                return jsonify({"error": "URL이 비어있습니다"}), 400

            source_label = f"URL: {url}"

            if enable_rendering:
                visible_text, hidden_text, iframe_text, screenshot_path = render_page_text_playwright(url)
            else:
                visible_text = url

            if enable_ocr and screenshot_path and api_key:
                img = Image.open(screenshot_path)
                ocr_text = ocr_google_vision(img, api_key)

        elif source_type == "TEXT":
            input_text = (request.form.get("input_text") or "").strip()
            if not input_text:
                return jsonify({"error": "입력 텍스트가 비어있습니다"}), 400

            source_label = "TEXT 입력"
            visible_text = normalize_text(input_text)

        elif source_type == "FILE":
            if "file" not in request.files:
                return jsonify({"error": "파일이 첨부되지 않았습니다"}), 400

            f = request.files["file"]
            if not f or not f.filename:
                return jsonify({"error": "파일명이 유효하지 않습니다"}), 400

            filename = safe_filename(f.filename)
            ext = get_ext(filename)
            save_path = os.path.join(UPLOAD_DIR, f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{filename}")
            f.save(save_path)

            source_label = f"FILE: {filename}"

            if ext in ALLOWED_PPT_EXT:
                return jsonify({"error": "PPT(.ppt)는 구버전 포맷으로 자동 텍스트 추출이 어렵습니다.\nPowerPoint에서 '다른 이름으로 저장' → PPTX로 저장 후 다시 업로드해주세요."}), 400

            if ext in ALLOWED_PPTX_EXT:
                return jsonify({"error": "PPTX는 먼저 분석(Preview) 후 실행됩니다. 파일 선택 시 자동 분석을 진행해주세요."}), 400

            elif ext in ALLOWED_IMAGE_EXT:
                img = Image.open(save_path)
                screenshot_path = save_path
                if enable_ocr and api_key:
                    ocr_text = ocr_google_vision(img, api_key)
                    if not ocr_text:
                        ocr_text = "(OCR 결과 없음)"
                visible_text = ""  # 이미지에는 텍스트 없음

            else:
                return jsonify({"error": f"지원하지 않는 파일 형식입니다: {ext}\n지원: PNG/JPG/WEBP, PPTX"}), 400

        elif source_type == "PPTX_TOKEN":
            token = (request.form.get("pptx_token") or "").strip()
            if not token or token not in PPTX_TOKENS:
                return jsonify({"error": "유효하지 않은 PPTX 토큰입니다. 다시 업로드/분석해주세요."}), 400

            pptx_ocr = request.form.get("pptx_ocr") == "1"

            info = PPTX_TOKENS[token]
            path = info["path"]
            filename = info["filename"]
            source_label = f"PPTX: {filename}"

            # 슬라이드 추출
            slides_data = pptx_extract_slides(path)

            # OCR(컨펌 후)
            if pptx_ocr:
                if not api_key:
                    return jsonify({"error": "PPTX OCR 실행을 위해 API Key가 필요합니다."}), 400
                slides_data = pptx_run_image_ocr(slides_data, api_key)

            # 합치기
            visible_text, ocr_text, slide_context_summary = pptx_merge_text(slides_data)

        else:
            return jsonify({"error": f"지원하지 않는 source_type: {source_type}"}), 400

        # -----------------------------
        # 2) 페이지 유형 분류
        # -----------------------------
        all_text = "\n".join([visible_text, hidden_text, iframe_text, ocr_text])
        page_type, page_type_reason = classify_page_type(all_text)

        # -----------------------------
        # 3) 규칙 기반 검사
        # -----------------------------
        score, findings, satisfied_hidden, rule_audit, passed_items, na_items, _ = run_rules(
            visible_text, hidden_text, iframe_text, ocr_text,
            enable_hidden=enable_hidden if source_type == "URL" else False,
            enable_ocr=enable_ocr,
            page_type=page_type
        )

        # -----------------------------
        # 4) Gemini 최종 검토 (선택)
        # -----------------------------
        gemini_result = None
        if enable_gemini and api_key:
            slide_context = ""
            if slide_context_summary:
                slide_context = f"\n{slide_context_summary}\n"
            gemini_result = review_with_gemini_rest(
                source_label, visible_text, hidden_text, iframe_text, ocr_text,
                rule_audit, findings, page_type, api_key, slide_context=slide_context
            )

        # -----------------------------
        # 5) 리포트 생성
        # -----------------------------
        options = {
            "source_type": source_type,
            "enable_rendering": enable_rendering,
            "enable_hidden": enable_hidden,
            "enable_ocr": enable_ocr,
            "enable_gemini": enable_gemini,
            "dev_mode": dev_mode,
            "pptx_ocr": request.form.get("pptx_ocr") == "1"
        }

        report_filename = generate_report_html_v23(
            source_label, score, findings, visible_text, hidden_text, iframe_text, ocr_text,
            screenshot_path, satisfied_hidden, rule_audit, passed_items, na_items,
            options, page_type, page_type_reason, gemini_result,
            slides_data=slides_data
        )

        return jsonify({"report": report_filename})

    except Exception as e:
        import traceback
        return jsonify({"error": f"{str(e)}\n\n{traceback.format_exc()}"}), 500


# ------------------------------------------------------------
# 리포트 생성(v2.3): 슬라이드별 표시 포함
# ------------------------------------------------------------
def generate_report_html_v23(
        source_label, score, findings, visible_text, hidden_text, iframe_text, ocr_text,
        screenshot_path, satisfied_hidden, rule_audit, passed_items, na_items,
        options, page_type, page_type_reason, gemini_result,
        slides_data=None
):
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    # Gemini 섹션
    gemini_section = ""
    if gemini_result:
        confidence = int(gemini_result.get("confidence_score", 0))
        confidence_color = "#1e8e3e" if confidence >= 80 else "#fbbc04" if confidence >= 60 else "#d11a2a"

        additional_findings_html = ""
        if gemini_result.get('additional_findings'):
            additional_findings_html = "<h3>추가 발견 사항</h3><ul>"
            for f in gemini_result['additional_findings']:
                sev = (f.get("심각도") or "LOW").upper()
                sev_class = f"badge-{sev.lower()}" if sev.lower() in ["critical", "high", "medium", "low"] else "badge-low"
                severity_badge = f"<span class='badge {sev_class}'>{sev}</span>"
                additional_findings_html += f"""
                <li style="margin-bottom:10px;">
                    {severity_badge} <strong>{escape_html(f.get('항목', ''))}</strong>
                    <br/>{escape_html(f.get('설명', ''))}
                    <br/><small style="color:#666;">근거: {escape_html(f.get('근거', ''))}</small>
                </li>
                """
            additional_findings_html += "</ul>"

        recommendations_html = ""
        if gemini_result.get('recommendations'):
            recommendations_html = "<h3>개선 권고 사항</h3><ul>"
            for r in gemini_result['recommendations']:
                recommendations_html += f"<li>{escape_html(r)}</li>"
            recommendations_html += "</ul>"

        raw_html = ""
        if options.get("dev_mode"):
            raw_json = gemini_result.get("raw_json")
            raw_text = gemini_result.get("raw_text") or ""

            if raw_json:
                raw_html += f"""
                <details style="margin-top: 15px;">
                  <summary style="cursor:pointer;font-weight:600;">Gemini 원문 JSON 보기 (개발자 모드)</summary>
                  <pre style="background:#fff;padding:12px;border-radius:6px;overflow:auto;max-height:380px;">
{escape_html(json.dumps(raw_json, ensure_ascii=False, indent=2))}
                  </pre>
                </details>
                """

            if raw_text:
                raw_html += f"""
                <details style="margin-top: 10px;">
                  <summary style="cursor:pointer;font-weight:600;">Gemini 원문 응답(raw_text) 보기 (개발자 모드)</summary>
                  <pre style="background:#fff;padding:12px;border-radius:6px;overflow:auto;max-height:380px;white-space:pre-wrap;">
{escape_html(raw_text)}
                  </pre>
                </details>
                """

        gemini_section = f"""
        <div class="section gemini-section">
            <h2>🤖 Gemini AI 최종 검토</h2>
            <div class="gemini-confidence">
                <div style="display: flex; align-items: center; justify-content: space-between;">
                    <span><strong>AI 신뢰도 점수:</strong></span>
                    <span style="font-size: 32px; font-weight: bold; color: {confidence_color};">
                        {confidence}/100
                    </span>
                </div>
            </div>
            <div class="gemini-review">
                <h3>종합 검토 의견</h3>
                <p style="line-height: 1.6;">{escape_html(gemini_result.get('gemini_review', ''))}</p>
            </div>
            {additional_findings_html}
            {recommendations_html}
            {raw_html}
        </div>
        """

    # Findings
    findings_html = ""
    if findings:
        for f in findings:
            sev = f.get("severity", "LOW")
            sev_class = {"CRITICAL": "badge-critical", "HIGH": "badge-high", "MEDIUM": "badge-medium",
                         "LOW": "badge-low"}.get(sev, "badge-low")
            findings_html += f"""
            <div class="finding">
                <span class="badge {sev_class}">{sev}</span>
                <strong>{escape_html(f.get('title', ''))}</strong>
                <br/>{escape_html(f.get('desc', ''))}
            </div>
            """
    else:
        findings_html = "<p style='color:#1e8e3e;'>✅ 위반 사항이 발견되지 않았습니다.</p>"

    # Rule audit
    rule_audit_html = "<table><tr><th>규칙</th><th>상태</th><th>검토 사유</th></tr>"
    for r in (rule_audit or [])[:15]:
        rule_audit_html += f"<tr><td>{escape_html(r.get('title', ''))}</td><td>{r.get('status', 'N/A')}</td><td>{escape_html(r.get('reason', ''))}</td></tr>"
    rule_audit_html += "</table>"

    # PPTX 슬라이드 섹션
    slides_section = ""
    if slides_data:
        blocks = []
        for s in slides_data:
            no = s["slide_no"]
            txt = s.get("text") or ""
            ocr = s.get("ocr") or ""
            img_count = len(s.get("images", []))

            blocks.append(f"""
            <details style="margin:10px 0;">
              <summary style="cursor:pointer;font-weight:700;">
                슬라이드 {no} (텍스트 {len(txt)}자 / OCR {len(ocr)}자 / 이미지 {img_count}개)
              </summary>
              <div style="padding:12px; background:#fff; border-radius:8px; margin-top:8px;">
                <h4 style="margin:0 0 8px 0;">슬라이드 텍스트</h4>
                <pre style="white-space:pre-wrap;">{escape_html(txt[:5000])}</pre>
                <h4 style="margin:10px 0 8px 0;">슬라이드 이미지 OCR</h4>
                <pre style="white-space:pre-wrap;">{escape_html(ocr[:5000])}</pre>
              </div>
            </details>
            """)

        slides_section = f"""
        <div class="section">
            <h2>📑 PPTX 슬라이드별 텍스트/이미지 OCR</h2>
            {"".join(blocks)}
        </div>
        """

    html = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Legal Check Report - {escape_html(source_label)}</title>
    <style>
        body {{ font-family: 'Segoe UI', sans-serif; margin: 0; padding: 20px; background: #f5f5f5; }}
        .container {{ max-width: 1200px; margin: 0 auto; background: white; padding: 30px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); }}
        h1 {{ color: #1a73e8; border-bottom: 3px solid #1a73e8; padding-bottom: 10px; }}
        h2 {{ color: #34a853; margin-top: 30px; }}
        h3 {{ color: #333; margin-top: 20px; }}
        .score {{ font-size: 48px; font-weight: bold; color: {'#d11a2a' if score < 70 else '#fbbc04' if score < 90 else '#1e8e3e'}; }}
        .section {{ margin: 20px 0; padding: 20px; background: #f8f9fa; border-radius: 8px; }}
        .gemini-section {{ background: #e8f0fe; border-left: 4px solid #1a73e8; }}
        .gemini-confidence {{ font-size: 18px; margin: 15px 0; padding: 15px; background: white; border-radius: 4px; }}
        .gemini-review {{ margin: 15px 0; padding: 15px; background: white; border-radius: 4px; }}
        .finding {{ margin: 10px 0; padding: 15px; background: white; border-left: 4px solid #d11a2a; border-radius: 4px; }}
        table {{ width: 100%; border-collapse: collapse; margin: 15px 0; background: white; }}
        th, td {{ padding: 12px; text-align: left; border-bottom: 1px solid #ddd; }}
        th {{ background: #f1f3f4; font-weight: 600; }}
        .badge {{ padding: 4px 8px; border-radius: 4px; font-size: 12px; font-weight: bold; color: white; }}
        .badge-critical {{ background: #d11a2a; }}
        .badge-high {{ background: #ea4335; }}
        .badge-medium {{ background: #fbbc04; color: #333; }}
        .badge-low {{ background: #34a853; }}
        .meta {{ color: #666; font-size: 14px; }}
        pre {{ background:#fff; padding:12px; border-radius:8px; overflow:auto; }}
    </style>
</head>
<body>
    <div class="container">
        <h1>법적 검토 리포트</h1>
        <p class="meta"><strong>검토 대상:</strong> {escape_html(source_label)}</p>
        <p class="meta"><strong>검사 시각:</strong> {now}</p>
        <p class="meta"><strong>입력 방식:</strong> {escape_html(options.get("source_type", ""))}</p>
        <p class="meta"><strong>PPTX OCR:</strong> {"ON" if options.get("pptx_ocr") else "OFF"}</p>
        <p class="meta"><strong>페이지 유형:</strong> {page_type}</p>
        <p class="meta"><strong>페이지 유형 근거:</strong> {" / ".join(page_type_reason or [])}</p>

        <div class="section">
            <h2>📊 규칙 기반 검사 점수</h2>
            <div class="score">{score}점</div>
            <p class="meta">100점 만점 기준, 규칙 위반 시 감점</p>
        </div>

        {gemini_section}

        <div class="section">
            <h2>⚠️ 위반 사항 ({len(findings)}건)</h2>
            {findings_html}
        </div>

        <div class="section">
            <h2>📋 규칙별 검사 결과 (상위 15개)</h2>
            {rule_audit_html}
        </div>

        {slides_section}

        <div class="section">
            <h2>📄 수집된 텍스트 (요약)</h2>
            <h3>Visible/Text</h3>
            <pre>{escape_html((visible_text or "")[:4000])}</pre>
            <h3>OCR</h3>
            <pre>{escape_html((ocr_text or "")[:2500])}</pre>
        </div>
    </div>
</body>
</html>
    """

    fname = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
    outpath = os.path.join(OUTPUT_DIR, fname)
    with open(outpath, "w", encoding="utf-8") as f:
        f.write(html)

    return fname


@app.route("/report/<path:filename>")
def serve_report(filename):
    return send_from_directory(OUTPUT_DIR, filename)


@app.route("/uploads/<path:filename>")
def serve_uploads(filename):
    return send_from_directory(UPLOAD_DIR, filename)


# ------------------------------------------------------------
# Main
# ------------------------------------------------------------
if __name__ == "__main__":
    # 클라우드 배포용 포트 설정
    port = int(os.environ.get("PORT", 8080))
    app.run(host="0.0.0.0", port=port, debug=False)